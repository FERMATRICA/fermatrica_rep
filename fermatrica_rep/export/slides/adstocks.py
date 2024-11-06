"""
Generate slide with adstock plot: adstock (decay, carry-over) effect for number of variables
(marketing tools).
"""


import pandas as pd

import lxml.etree as etree
from pptx.presentation import Presentation
from pptx.chart.data import ChartData
from pptx.util import Inches, Pt, Cm
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION, XL_DATA_LABEL_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx import enum as pptx_enum

from fermatrica import Model

from fermatrica_rep.model_rep import ModelRep
from fermatrica_rep.export.basics import set_chart_colors_fill
from fermatrica_rep.adstocks import adstocks_data as adstocks_data
from fermatrica_rep.options.define import OptionSettings


def create(prs: Presentation,
           model: "Model | list",
           ds: pd.DataFrame | list,
           model_rep,  #: "ModelRep"
           cln_meas: list | tuple,
           option_dict: dict | list = None,
           option_settings: "OptionSettings" = None,
           cln_dim: list | tuple = ('superbrand', 'master', 'bs_key', 'date', 'listed', 'kpi_coef'),
           n: int = 50
           ):
    """
    Create adstock slide and add to `prs` PPTX presentation.

    :param prs: Presentation object from python_pptx package
    :param model: Model object or list of Model objects
    :param ds: dataset or list of datasets
    :param model_rep: ModelRep object (export settings) of list of ModelRep objects
    :param option_dict: budget option / scenario to calculate as dictionary or list of option dictionaries
    :param option_settings: OptionSettings object (option setting: target period etc.) or list of OptionSettings
    :param cln_meas: column names to be used as measurements
    :param cln_dim: column names to be used as dimensions
    :param n: number of observations per column
    :return: Presentation object from python_pptx package
    """
    if isinstance(model, list):
        language = model_rep[-1].language
        superbrand = model[-1].conf.target_superbrand
    else:
        language = model_rep.language
        superbrand = model.conf.target_superbrand

    # data prep
    adstocks_dataframe = adstocks_data(model=model
                                       , ds=ds
                                       , model_rep=model_rep
                                       , superbrand=superbrand
                                       , option_settings=option_settings
                                       , option_dict=option_dict
                                       , cln_dim=cln_dim
                                       , cln_meas=cln_meas
                                       , n=n
                                       )

    superbrand = superbrand.title()

    # prepare slide in presentation

    if isinstance(model_rep, list):
        model_rep = model_rep[-1]
        model = model[-1]

    slide_layout = model_rep.pptx_cnf['Blank_slide']

    slide_height = model_rep.pptx_cnf['slide_height']
    slide_width = model_rep.pptx_cnf['slide_width']

    slide = prs.slides.add_slide(slide_layout)

    slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'adstocks') &
                                         (model_rep.vis_dict['variable'] == 'slide_title'), language].iloc[0]

    slide_title = slide_title.replace('***id', superbrand)

    slide.shapes[0].text_frame.text = slide_title
    slide.shapes[0].text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_gross']
    slide.shapes[0].text_frame.paragraphs[0].font.name = model_rep.pptx_cnf['font_family_header']

    # mock text
    x, y = 0.55 * slide_width, slide.shapes[0].height + slide.shapes[0].top * 1.7
    cx, cy = 0.35 * slide_width, 0.55 * slide_height

    textbox = slide.shapes.add_textbox(x, y, cx, cy)
    textbox.text = model_rep.pptx_cnf['mock_text']

    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.font.size = model_rep.pptx_cnf["font_size_main"]
    paragraph.font.name = model_rep.pptx_cnf["font_family_body"]
    textbox.text_frame.word_wrap = True

    if len(adstocks_dataframe) == 0:
        return prs

    adstocks_dataframe = adstocks_dataframe * 100
    adstocks_dataframe[adstocks_dataframe < 1e-4] = "#N/A"

    for i in adstocks_dataframe.columns:
        disp_var = model.conf.model_rhs.loc[(model.conf.model_rhs.if_active == 1) & (model.conf.model_rhs.token.str.contains(i)), "display_var"].iloc[0]
        if disp_var in model_rep.palette_tools:
            adstocks_dataframe.rename(columns={i : disp_var}, inplace=True)

    # building chart
    chart_data_val = ChartData()
    chart_data_val.categories = pd.Series(adstocks_dataframe.index)

    # del adstocks_dataframe["date"]

    for i, j in adstocks_dataframe.items():
        chart_data_val.add_series(name=i,
                                  values=j)

    x, y = 0.05 * slide_width, slide.shapes[0].height + slide.shapes[0].top
    cx, cy = 0.45 * slide_width, 0.65 * slide_height

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_val
    )

    chart = graphic_frame.chart
    chart = set_chart_colors_fill(chart, model_rep, RGBColor(170, 170, 170))

    for sr in chart.plots[0].series:
        sr.format.line.width = Pt(2.5)
        sr.smooth = True

    etree.SubElement(chart._chartSpace.xpath("//c:valAx")[0], '{http://schemas.openxmlformats.org/drawingml/2006/chart}crossBetween').attrib['val'] = 'midCat'

    skip = str(round(n / 5))

    if int(skip) < 1:
        skip = "1"

    etree.SubElement(chart._chartSpace.xpath("//c:catAx")[0], '{http://schemas.openxmlformats.org/drawingml/2006/chart}tickLblSkip').attrib['val'] = skip
    etree.SubElement(chart._chartSpace.xpath("//c:catAx")[0], '{http://schemas.openxmlformats.org/drawingml/2006/chart}tickMarkSkip').attrib['val'] = skip

    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.chart_title.text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_main']

    category_axis = chart.category_axis
    category_axis.has_title = True
    category_axis_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "adstocks") & (
                model_rep.vis_dict['variable'] == "axis_x"), model_rep.language].iloc[0]

    category_axis_title = category_axis_title.replace('***period', model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "adstocks") & (
                model_rep.vis_dict['variable'] == model.conf.period_var), model_rep.language].iloc[0])

    category_axis.axis_title.text_frame.text = category_axis_title
    category_axis.axis_title.text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_small']

    value_axis = chart.value_axis
    value_axis.has_title = True
    value_axis_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "adstocks") & (
                model_rep.vis_dict['variable'] == "axis_y"), model_rep.language].iloc[0]

    value_axis.axis_title.text_frame.text = value_axis_title
    value_axis.axis_title.text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_small']

    value_axis.minimum_scale = 0
    value_axis.maximum_scale = 100
    value_axis.major_unit = 20

    value_axis.tick_labels.number_format = "0\%"

    value_axis.has_minor_gridlines = False
    value_axis.has_major_gridlines = True

    chart.font.size = model_rep.pptx_cnf['font_size_small']
    chart.font.name = model_rep.pptx_cnf['font_family_body']

    chart.has_legend = True
    chart.legend.include_in_layout = False

    text_frame = chart.chart_title.text_frame
    text_frame.text = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'adstocks') &
                                             (model_rep.vis_dict['variable'] == 'chart_title'), language].iloc[0]
    text_frame.text = text_frame.text.replace('***id', superbrand)

    text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_main']

    return prs

