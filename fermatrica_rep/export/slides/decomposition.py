"""
Generate slide with decomposition plot. Could be used both in retro analysis and in option reporting,
for superbrand as a whole and for specific number of variables known as "SKU"
"""


import pandas as pd
import copy

from pptx.presentation import Presentation
from pptx.chart.data import ChartData
from pptx.util import Inches, Pt, Cm
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor
from pptx import enum as pptx_enum

from fermatrica_utils import date_to_period, groupby_eff, select_eff

from fermatrica.model.model import Model

from fermatrica_rep.export.basics import set_chart_colors_fill
from fermatrica_rep.model_rep import ModelRep


def create(prs: Presentation,
           model: "Model",
           model_rep: "ModelRep",
           split_m_m: pd.DataFrame,
           option_name: str,
           brands: list,
           period: str = 'day',
           group_var: list | tuple = ('superbrand', 'market'),
           plot_type: str = 'brand',
           bs_key_filter: list | tuple | None = None,
           if_volume: bool = True,
           show_future: bool | str = True,
           contour_line: bool = True
           ):
    """
    Create (dynamic) decomposition slide and add to `prs` PPTX presentation.

    :param prs: Presentation object from python_pptx package
    :param model_rep: ModelRep object (reporting settings)
    :param model: Model object
    :param split_m_m: prepared dataset (see `fermatrica_rep.extract_effect()`)
    :param option_name: name of the option to be reported. To be used only as title here, no impact on program behaviour
    :param brands: superbrand string names as list
    :param period: time period / interval to group by: 'day', 'week', 'month', 'quarter', 'year'
    :param group_var: group entities by variables (list or tuple of strings)
    :param plot_type: 'retro', 'brand' or 'sku'. Used for naming only
    :param bs_key_filter: list or tuple of 'bs_key' values to preserve
    :param if_volume: optimize volume or value KPI
    :param show_future: show future periods or not
    :param contour_line: add contours or not
    :return: Presentation object from python_pptx package
    """

    language = model_rep.language

    if hasattr(model.conf, 'display_units'):
        units = model.conf.display_units
    else:
        units = 'packs'

    # prepare slide in presentation

    slide_layout = model_rep.pptx_cnf['Blank_slide']

    slide_height = model_rep.pptx_cnf['slide_height']
    slide_width = model_rep.pptx_cnf['slide_width']

    slide = prs.slides.add_slide(slide_layout)

    if ('retro' in plot_type) & (if_volume is True):
        slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'factor_decomposition') &
                                             (model_rep.vis_dict['variable'] == 'title_vol'), language].iloc[0]
    elif ('retro' in plot_type) & (if_volume is False):
        slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'factor_decomposition') &
                                             (model_rep.vis_dict['variable'] == 'title_val'), language].iloc[0]
    else:
        match (plot_type, if_volume):
            case ('brand', True):
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.superbrand') &
                                                     (model_rep.vis_dict['variable'] == 'title'), language].iloc[0]
            case ('brand', False):
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.superbrand') &
                                                     (model_rep.vis_dict['variable'] == 'title_value'), language].iloc[0]
            case ('sku', True):
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.sku') &
                                                     (model_rep.vis_dict['variable'] == 'title'), language].iloc[0]
            case ('sku', False):
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.sku') &
                                                     (model_rep.vis_dict['variable'] == 'title_value'), language].iloc[0]
            case _:
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.superbrand') &
                                                     (model_rep.vis_dict['variable'] == 'title'), language].iloc[0]

    slide_title = slide_title.replace('***id', option_name)
    slide_title = slide_title.replace('***units', model_rep.vis_dict.loc[
        (model_rep.vis_dict['section'] == "shared") & (model_rep.vis_dict['variable'] == units), language].iloc[0])

    slide.shapes[0].text_frame.text = slide_title
    slide.shapes[0].text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_gross']
    slide.shapes[0].text_frame.paragraphs[0].font.name = model_rep.pptx_cnf['font_family_header']

    # data prep for decompose
    split = copy.deepcopy(split_m_m)
    split = split[split['bs_key'].isin(bs_key_filter)]

    df_decompose, sku = _decompose_chart_data(split,
                                              brands=brands,
                                              show_future=show_future,
                                              period=period,
                                              group_var=group_var,
                                              if_volume=if_volume)

    # constructing graphic
    chart_data = ChartData()
    chart_data.categories = df_decompose.index
    for col_name, col in df_decompose.items():
        chart_data.add_series(col_name, col)

    x, y = (0.5 - 0.45) * slide_width, slide.shapes[0].top + slide.shapes[0].height
    cx, cy = 0.9 * slide_width, 0.7 * slide_height

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    chart = set_chart_colors_fill(chart, model_rep)

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.LEFT
    chart.legend.include_in_layout = False
    chart.legend.font.name = model_rep.pptx_cnf['font_family_body']
    chart.legend.font.size = model_rep.pptx_cnf['font_size_small']

    chart.plots[0].gap_width = 0

    chart.font.size = model_rep.pptx_cnf['font_size_small']
    chart.font.name = model_rep.pptx_cnf['font_family_body']

    for x in chart.plots[0].series:
        x.invert_if_negative = False
        if contour_line:
            x.format.line.width = Pt(0.5)
            x.format.line.color.rgb = RGBColor(200, 200, 200)

    category_axis = chart.category_axis
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    category_axis.has_minor_gridlines = True
    category_axis.has_major_gridlines = False

    if period == 'month':
        chart._chartSpace.xpath("//c:dateAx")[0].xpath("//c:baseTimeUnit")[0].attrib['val'] = 'months'
    elif period == 'year':
        chart._chartSpace.xpath("//c:dateAx")[0].xpath("//c:baseTimeUnit")[0].attrib['val'] = 'years'

    value_axis = chart.value_axis
    value_axis.axis_title.text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_main']

    chart.has_title = True
    text_frame = chart.chart_title.text_frame
    if if_volume:
        text_frame.text = f"{sku}: {model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'plot_efficiency') & (model_rep.vis_dict['variable'] == 'factors_title'), language].iloc[0]}"
    else:
        text_frame.text = f"{sku}: {model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'plot_efficiency') & (model_rep.vis_dict['variable'] == 'factors_title_val'), language].iloc[0]}"

    text_frame.text = text_frame.text.replace('***units', model_rep.vis_dict.loc[
        (model_rep.vis_dict['section'] == "shared") & (model_rep.vis_dict['variable'] == units), language].iloc[0])
    text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_main']

    return prs


def _decompose_chart_data(split: pd.DataFrame,
                          brands: list,
                          show_future: bool | str = True,
                          period: str = 'day',
                          group_var: list | tuple = ('superbrand', 'market'),
                          if_volume: bool = True) -> tuple:
    """
    Convert 'common' decomposed dataset to dataset prepared to be plotted.

    :param split: prepared dataset (see `fermatrica_rep.extract_effect()`) filtered by bs_key
    :param brands: superbrand string names as list
    :param period: time period / interval to group by: 'day', 'week', 'month', 'quarter', 'year'
    :param group_var: group entities by variables (list or tuple of strings)
    :param if_volume: optimize volume or value KPI
    :param show_future: show future periods or not
    :return: tuple of pandas DataFrame and sku name as string
    """

    if not brands:
        brands = split['superbrand'].unique()

    if show_future == 'True':
        show_future = True
    elif show_future == 'False':
        show_future = False

    if period in ['day', 'd', 'D']:
        split['date'] = split['date'].dt.floor(freq='D')
    elif period in ['week', 'w', 'W']:
        split['date'] = split['date'].dt.to_period('W').dt.start_time
    elif period in ['month', 'm', 'M']:
        split['date'] = split['date'].dt.to_period('M').dt.start_time
    elif period in ['quarter', 'q', 'Q']:
        split['date'] = split['date'].dt.to_period('Q').dt.start_time
    elif period in ['year', 'y', 'Y']:
        split['date'] = split['date'].dt.to_period('Y').dt.start_time

    if show_future:
        mask = split['superbrand'].isin(brands) & split['listed'].isin([2, 3, 4])
    else:
        mask = split['superbrand'].isin(brands) & split['listed'].isin([2, 3])

    if if_volume is False:
        split['value'] = split['value_rub']

    tmp = groupby_eff(split, list(set(['date'] + group_var + ['variable'])), ['value'], mask, sort=False)
    tmp = tmp.value.sum().reset_index()
    tmp['sku'] = ''

    for i in group_var:
        tmp['sku'] = tmp['sku'] + ' ' + tmp[i].astype('str').str.title()
    tmp['sku'] = tmp['sku'].str.lstrip()

    sku = tmp['sku'].unique()[0]

    tmp_wide = pd.pivot(tmp, index='date', columns='variable', values='value')

    # delete all columns with only 0 values
    tmp_wide = tmp_wide.loc[:, (tmp_wide != 0).any(axis=0)]

    return tmp_wide, sku
