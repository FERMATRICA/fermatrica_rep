"""
Generate slide with waterfall (static decomposition) plot. Could be used both in retro analysis and in option
reporting. As for now only superbrand export is available (no SKU's). `bs_key_filter` could also be used
to filter only some of bs_key of the brands to be reported.
"""


import pandas as pd
import numpy as np
import copy

from pptx.presentation import Presentation
from pptx.chart.data import ChartData
from pptx.util import Inches, Pt, Cm
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION, \
    XL_DATA_LABEL_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx import enum as pptx_enum

from fermatrica_utils import groupby_eff, hex_to_rgb

from fermatrica.model.model import Model
from fermatrica_rep.model_rep import ModelRep


def create(prs: Presentation,
           model: "Model",
           model_rep: "ModelRep",
           split_m_m: pd.DataFrame,
           option_name: str,
           brands: list,
           date_start: str,
           date_end: str,
           absolute_sort: str | bool = False,
           plot_type: str = 'brand',
           bs_key_filter: list | tuple | None = None,
           if_volume: bool = True
           ):
    """
    Create  waterfall (static decomposition) slide and add to `prs` PPTX presentation.
    As for now only superbrand export is available (no SKU's). `bs_key_filter` could also be used
    to filter only some of bs_key of the brands to be reported.

    :param prs: Presentation object from python_pptx package
    :param model_rep: ModelRep object (reporting settings)
    :param model: Model object
    :param split_m_m: prepared dataset (see `fermatrica_rep.extract_effect()`)
    :param option_name: name of the option to be reported. To be used only as title here, no impact on program behaviour
    :param brands: superbrand string names as list
    :param date_start: start of the period
    :param date_end: end of the period
    :param absolute_sort: sort by absolute values or take into account sign
    :param plot_type: 'retro', 'brand' or 'sku'. Used for naming only
    :param bs_key_filter: list or tuple of 'bs_key' values to preserve
    :param if_volume: optimize volume or value KPI
    :return: Presentation object from python_pptx package
    """

    language = model_rep.language

    if hasattr(model.conf, 'display_units'):
        units = model.conf.display_units
    else:
        units = 'packs'

    # prepare slide in presentation

    slide_layout = model_rep.pptx_cnf['Blank_slide']

    slide_height = model_rep.pptx_cnf['slide_height']
    slide_width = model_rep.pptx_cnf['slide_width']

    slide = prs.slides.add_slide(slide_layout)

    if ('retro' in plot_type) & (if_volume is True):
        slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'factor_decomposition') &
                                             (model_rep.vis_dict['variable'] == 'title_vol'), language].iloc[0]
    elif ('retro' in plot_type) & (if_volume is False):
        slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'factor_decomposition') &
                                             (model_rep.vis_dict['variable'] == 'title_val'), language].iloc[0]
    else:
        match (plot_type, if_volume):
            case ('brand', True):
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.superbrand') &
                                                     (model_rep.vis_dict['variable'] == 'title'), language].iloc[0]
            case ('brand', False):
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.superbrand') &
                                                     (model_rep.vis_dict['variable'] == 'title_value'), language].iloc[0]
            case ('sku', True):
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.sku') &
                                                     (model_rep.vis_dict['variable'] == 'title'), language].iloc[0]
            case ('sku', False):
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.sku') &
                                                     (model_rep.vis_dict['variable'] == 'title_value'), language].iloc[0]
            case _:
                slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.effects.superbrand') &
                                                     (model_rep.vis_dict['variable'] == 'title'), language].iloc[0]

    slide_title = slide_title.replace('***id', option_name)
    slide_title = slide_title.replace('***units', model_rep.vis_dict.loc[
        (model_rep.vis_dict['section'] == "shared") & (model_rep.vis_dict['variable'] == units), language].iloc[0])

    slide.shapes[0].text_frame.text = slide_title
    slide.shapes[0].text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_gross']
    slide.shapes[0].text_frame.paragraphs[0].font.name = model_rep.pptx_cnf['font_family_header']

    # mock text
    x, y = 0.05 * slide_width, slide.shapes[0].height + slide.shapes[0].top
    cx, cy = 0.25 * slide_width, 0.7 * slide_height
    textbox = slide.shapes.add_textbox(x, y, cx, cy)
    textbox.text = model_rep.pptx_cnf['mock_text']

    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.font.size = model_rep.pptx_cnf["font_size_main"]
    paragraph.font.name = model_rep.pptx_cnf["font_family_body"]
    textbox.text_frame.word_wrap = True

    # data prep for waterfall
    split = copy.deepcopy(split_m_m)

    if bs_key_filter is not None:
        split = split[split['bs_key'].isin(bs_key_filter)]

    waterfall_data = _waterfall_chart_data(split=split,
                                           brands=brands,
                                           date_start=date_start,
                                           date_end=date_end,
                                           absolute_sort=absolute_sort,
                                           if_volume=if_volume)

    if len(waterfall_data) == 0:
        return prs

    # constructing graphic
    chart_data = ChartData()
    categories = np.where(waterfall_data['variable'].str.len() > 15,
                          waterfall_data['variable'].str[0:15] + '...',
                          waterfall_data['variable'])
    chart_data.categories = categories

    chart_data.add_series('base', waterfall_data['base'])
    chart_data.add_series('value', waterfall_data['value'])

    x, y = (0.5 - 0.45 + 0.25) * slide_width, slide.shapes[0].top + slide.shapes[0].height
    cx, cy = (0.9 - 0.25) * slide_width, 0.7 * slide_height

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    chart.plots[0].gap_width = 20

    chart_fill = chart.series[0].format.fill
    chart_fill.solid()
    chart_fill.fore_color.rgb = RGBColor(255, 0, 0)

    # ---add an `a:alpha` child element---
    solidFill = chart_fill.fore_color._xFill
    alpha = OxmlElement('a:alpha')
    alpha.set('val', '0')
    solidFill.srgbClr.append(alpha)

    chart.font.size = model_rep.pptx_cnf['font_size_small']
    chart.font.name = model_rep.pptx_cnf['font_family_body']

    chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.HIGH

    value_axis = chart.value_axis
    value_axis.has_title = True

    if language == 'english':
        value_axis_title = 'Factors contribution to forecast, %'
    else:
        value_axis_title = 'Вклад факторов в прогноз, %'

    value_axis.axis_title.text_frame.text = value_axis_title
    value_axis.axis_title.text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_small']

    for idx, point in enumerate(chart.series[1].points):

        point.data_label.has_text_frame = True
        point.data_label.text_frame.text = waterfall_data['ratio_lbl'].iloc[idx]

        # point.data_label.font.bold = True
        # point.data_label.font.name = model_rep.pptx_cnf['font_family_body']
        # point.data_label.font.size = model_rep.pptx_cnf['font_size_footnote']

        point.data_label.position = XL_DATA_LABEL_POSITION.INSIDE_END

        if waterfall_data['variable'].iloc[idx] in model_rep.palette_tools:
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*hex_to_rgb(model_rep.palette_tools[waterfall_data['variable'].iloc[idx]]))

    chart.series[1].data_labels.font.bold = True
    chart.series[1].data_labels.font.size = model_rep.pptx_cnf['font_size_small']
    chart.series[1].data_labels.font.name = model_rep.pptx_cnf['font_family_body']

    return prs


def _waterfall_chart_data(split: pd.DataFrame,
                          brands: list,
                          date_start: str,
                          date_end: str,
                          absolute_sort: str | bool = False,
                          if_volume: bool = True):
    """

    Convert 'common' decomposed dataset to dataset prepared to be plotted.

    :param split: prepared dataset (see `fermatrica_rep.extract_effect()`) filtered by bs_key
    :param brands: superbrand string names as list
    :param date_start: start of the period
    :param date_end: end of the period
    :param absolute_sort: sort by absolute values or take into account sign
    :param if_volume: optimize volume or value KPI
    :return: pandas DataFrame prepared to be plotted
    """

    if not brands:
        brands = split['superbrand'].unique()

    if absolute_sort == 'True':
        absolute_sort = True
    elif absolute_sort == 'False':
        absolute_sort = False

    mask = (split['date'] >= pd.to_datetime(date_start)) & \
           (split['date'] <= pd.to_datetime(date_end)) & \
           (split['superbrand'].isin(brands))

    if if_volume is False:
        split['value'] = split['value_rub']

    split = groupby_eff(split, ['variable'], ['value'], mask, sort=False)['value'].sum().reset_index()
    split['abs'] = split['value'].abs()

    if absolute_sort:
        split.sort_values(by='abs', ascending=False, inplace=True)
    else:
        split.sort_values(by='value', ascending=False, inplace=True)

    split['ratio'] = (split['value'] / split['value'].sum()) * 100
    split['ratio_lbl'] = split['ratio'].round(2).astype(str) + '%'

    split['value_cumsum'] = split['value'].cumsum()
    split['ratio_cumsum'] = (split['value_cumsum'] / split['value'].sum()) * 100

    split = split[split['ratio'] != 0]

    split['base'] = np.where(split['ratio'] > 0, split['ratio_cumsum'] - split['ratio'], split['ratio_cumsum'])

    split['value'] = split['ratio'].abs()

    split = split.reset_index()

    return split[['variable', 'value', 'base', 'ratio_lbl']]
