"""
Generate 2 slides with efficiency curves (i.e. incremental KPI, profit, ROI) via full curves approach
(row of options to be calculated): short and long term effects.

To get stable results run `fermatrica.options.calc.option_report()`
with some fixed option (as 'zero') and `exact=False` before calling `fermatrica_rep.export.export_pptx.export_curves()`.
Period to take into account could be set via OptionSettings object.

Beware! Multiprocessing calculation to be used by functions of this file.
"""
import copy
from typing import Callable
import pandas as pd

from pptx.presentation import Presentation
from pptx.chart.data import ChartData
from pptx.util import Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
import lxml.etree as etree

from fermatrica_utils import hex_to_rgb

from fermatrica.model.model import Model

from fermatrica_rep.export.basics import set_chart_colors_line, set_chart_dashes
from fermatrica_rep.model_rep import ModelRep
from fermatrica_rep.options.define import OptionSettings
from fermatrica_rep.options.calc_multi import option_report_multi_post


def create(prs: Presentation,
           model: "Model | list",
           model_rep: "ModelRep | list",
           ds: pd.DataFrame | list,
           opt_set_crv,
           translation: dict | None = None,
           budget_step: int | float = 5,
           bdg_max: int | float = 301,
           fixed_vars: dict | None = {'price': 1},
           adhoc_curves_max_costs: "None | Callable" = None,
           if_exact: bool = True,
           cores: int = 4):
    """
    Generate efficiency curves slides (i.e. incremental KPI, profit, ROI) via full curves approach
    (row of options to be calculated): short and long term effects. 2 slides should be generated
    with short and long term curves respectively. 

    :param prs: Presentation object from python_pptx package
    :param model: Model object or list of Model objects
    :param model_rep: ModelRep object (export settings) or list of ModelRep objects
    :param ds: dataset or list of datasets
    :param opt_set_crv: OptionSetting object containing calculate settings
    :param translation: translation dict (from files like `options.xlsx`, `translation` sheet) or None,
        (i.e. defaults to trans_dict attribute of `model_rep` argument)
    :param budget_step: budget iteration step in millions, defaults to 1 (i.e. 1M)
    :param bdg_max: maximum budget size (all options with larger budgets to be dropped)
    :param adhoc_curves_max_costs: adhoc function to set maximum observed values for every variable (optional)
    :param fixed_vars: translation variables with their values to be fixed across grid
    :param if_exact: apply only to the specific time period, without next years
    :param cores: number of processor cores to use in calculations; None sets to all computer logical cores - 1
    :return: Presentation object from python_pptx package
    """

    if isinstance(model_rep, list):
        model_rep_main = model_rep[-1]
        date_max = ds[-1]['date'].dt.year.max()
    else:
        model_rep_main = model_rep
        date_max = ds['date'].dt.year.max()

    language = model_rep_main.language
    slide_width = model_rep_main.pptx_cnf['slide_width']
    slide_height = model_rep_main.pptx_cnf['slide_height']

    # fix for evading problem with pickling python-pptx object

    if isinstance(model_rep, list):
        tmp = [None] * len(model_rep)
        for i, mdr in enumerate(model_rep):
            tmp[i] = model_rep[i].pptx_cnf['Blank_slide']
            model_rep[i].pptx_cnf['Blank_slide'] = 0
    else:
        tmp = model_rep.pptx_cnf['Blank_slide']
        model_rep.pptx_cnf['Blank_slide'] = 0

    # load data
    curves_full_data = option_report_multi_post(model=model,
                                                ds=ds,
                                                model_rep=model_rep,
                                                opt_set=opt_set_crv,
                                                translation=translation,
                                                adhoc_curves_max_costs=adhoc_curves_max_costs,
                                                budget_step=budget_step,
                                                bdg_max=bdg_max,
                                                fixed_vars=fixed_vars,
                                                if_exact=if_exact,
                                                cores=cores)

    if isinstance(curves_full_data, list):
        curves_full_data = curves_full_data[-1]

    if isinstance(model_rep, list):

        for i, mdr in enumerate(model_rep):
            model_rep[i].pptx_cnf['Blank_slide'] = tmp[i]
    else:
        model_rep.pptx_cnf['Blank_slide'] = tmp

    # prepare slide in presentation

    slide_layout = model_rep_main.pptx_cnf['Blank_slide']

    # prepare colors
    tmp = curves_full_data.groupby('option').agg({'pred_long_val': 'nunique'})
    lst = tmp.loc[tmp["pred_long_val"] > 1].reset_index().option.tolist()
    lst.sort()

    model_rep_main.fill_colours_tools(lst)

    for i in lst:
        model_rep_main.palette_tools[i + "_extrapolated"] = model_rep_main.palette_tools[i]

    # ----------------- short effect ------------------

    slide = prs.slides.add_slide(slide_layout)

    slide = _slide_vis_worker(slide,
                              model_rep=model_rep_main,
                              opt_set_crv=opt_set_crv,
                              curves_full_data=curves_full_data,
                              plot_type='short',
                              year_end=date_max,
                              budget_step=budget_step,
                              bdg_max=bdg_max
                              )

    # ------------------ long effect ------------------

    slide = prs.slides.add_slide(slide_layout)

    slide = _slide_vis_worker(slide,
                              model_rep=model_rep_main,
                              opt_set_crv=opt_set_crv,
                              curves_full_data=curves_full_data,
                              plot_type='long',
                              year_end=date_max,
                              budget_step=budget_step,
                              bdg_max=bdg_max,
                              if_exact=if_exact)

    return prs


def _slide_vis_worker(slide,
                      model_rep: "ModelRep",
                      opt_set_crv: "OptionSettings",
                      curves_full_data: pd.DataFrame,
                      plot_type: str,
                      year_end: int,
                      budget_step: int | float = 5,
                      bdg_max: int | float = 301,
                      if_exact: bool = True
                      ):
    """
    Create one after another of two slides for `fermatrica_rep.export.slides.incr_roi.create()`.
    Calculations are already done, so this function works as visualiser only.
    
    This function calls minor workers to fill current file.

    :param slide: Slide object (from python_pptx package) to be filled
    :param model_rep: ModelRep object (reporting settings)
    :param opt_set_crv: OptionSetting object containing calculate settings
    :param curves_full_data:
    :param plot_type: 'short' or 'long'
    :param year_end: final year of the calculations; to be used in naming only
    :param budget_step: budget iteration step in millions, defaults to 1 (i.e. 1M)
    :param bdg_max: maximum budget size (all options with larger budgets to be dropped)
    :param if_exact: apply only to the specific time period, without next years
    :return:
    """

    if plot_type == 'long':
        n_years = year_end - opt_set_crv.date_start.year + 1
        axis_years = str(opt_set_crv.date_start.year) + '-' + str(year_end)
        slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'efficiency.3.year') &
                                             (model_rep.vis_dict['variable'] == 'title'), model_rep.language].iloc[0]
    else:
        n_years = 1
        axis_years = str(opt_set_crv.date_start.year)
        slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'efficiency.1.year') &
                                             (model_rep.vis_dict['variable'] == 'title'), model_rep.language].iloc[0]

    slide_height = model_rep.pptx_cnf['slide_height']
    slide_width = model_rep.pptx_cnf['slide_width']

    slide_title = slide_title.replace("***", axis_years)
    slide.shapes[0].text = slide_title
    paragraph = slide.shapes[0].text_frame.paragraphs[0]
    paragraph.font.size = model_rep.pptx_cnf["font_size_gross"]
    paragraph.font.name = model_rep.pptx_cnf["font_family_header"]
    slide.shapes[0].text_frame.word_wrap = True

    x, y = slide.shapes[0].left, slide.shapes[0].height + slide.shapes[0].top + 0.05 * slide_height
    cx, cy = slide.shapes[0].width, 0.2 * slide_height
    textbox = slide.shapes.add_textbox(x, y, cx, cy)
    textbox.text = model_rep.pptx_cnf['mock_text']

    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.font.size = model_rep.pptx_cnf["font_size_main"]
    paragraph.font.name = model_rep.pptx_cnf["font_family_body"]
    textbox.text_frame.word_wrap = True

    # Incremental sales
    chart_data = _chart_data_worker(curves_full_data, f'increment_value_{plot_type}', if_exact, n_years)

    x, y = (0.5 - 0.45) * slide_width, slide.shapes[1].top + slide.shapes[1].height
    cx, cy = 0.25 * slide_width, 0.5 * slide_height

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    chart = set_chart_dashes(chart, '_extrapolated')

    chart.has_legend = False
    chart = _set_chart_axis(chart, model_rep,
                            var='incr_sales_axis_y',
                            year=axis_years,
                            plot_type=plot_type,
                            bdg_max=bdg_max,
                            if_exact=if_exact,
                            budget_step=budget_step)

    chart = set_chart_colors_line(chart, model_rep, RGBColor(170, 170, 170))

    chart.font.size = model_rep.pptx_cnf['font_size_small']
    chart.font.name = model_rep.pptx_cnf['font_family_body']

    for sr in chart.plots[0].series:
        sr.format.line.width = Pt(2.5)
        sr.smooth = True

    if curves_full_data["pred_exact_val"][1] > 1e6:
        xml_node = etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/chart}dispUnits")
        etree.SubElement(xml_node, '{http://schemas.openxmlformats.org/drawingml/2006/chart}builtInUnit').attrib[
            'val'] = 'millions'
        dispUnitsLbl = etree.SubElement(xml_node,
                                        '{http://schemas.openxmlformats.org/drawingml/2006/chart}dispUnitsLbl')
        layout = etree.SubElement(dispUnitsLbl, '{http://schemas.openxmlformats.org/drawingml/2006/chart}layout')

        manual = etree.SubElement(layout, '{http://schemas.openxmlformats.org/drawingml/2006/chart}manualLayout')
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}xMode').attrib['val'] = 'edge'
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}yMode').attrib['val'] = 'edge'
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}x').attrib['val'] = '0.13'
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}y').attrib['val'] = '0.03'

        if model_rep.language == 'russian':
            tx = etree.SubElement(dispUnitsLbl, '{http://schemas.openxmlformats.org/drawingml/2006/chart}tx')
            rich = etree.SubElement(tx, '{http://schemas.openxmlformats.org/drawingml/2006/chart}rich')
            etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')

            p = etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')

            pPr = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')

            r = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
            rpr = etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')

            rpr.attrib['lang'] = 'ru-RU'
            rpr.attrib['dirty'] = '0'

            etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t').text = 'млн'

            etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}endParaRPr').attrib['lang'] = 'en-US'
            # etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}endParaRPr').attrib['dirty'] = '0'

        chart._chartSpace.xpath("//c:valAx")[0].insert(-1, xml_node)

    # profit
    chart_data = _chart_data_worker(curves_full_data, f'profit_{plot_type}', if_exact, n_years)

    x, y = (0.5 - 0.45 + 0.25) * slide_width, slide.shapes[1].top + slide.shapes[1].height
    cx, cy = 0.25 * slide_width, 0.5 * slide_height

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    chart = set_chart_dashes(chart, '_extrapolated')

    chart.has_legend = False

    chart = _set_chart_axis(chart
                            , model_rep,
                            var='profit_axis_y',
                            year=axis_years,
                            if_exact=if_exact,
                            plot_type=plot_type,
                            bdg_max=bdg_max,
                            budget_step=budget_step)

    chart = set_chart_colors_line(chart, model_rep, RGBColor(170, 170, 170))

    chart.font.size = model_rep.pptx_cnf['font_size_small']
    chart.font.name = model_rep.pptx_cnf['font_family_body']

    for sr in chart.plots[0].series:
        sr.format.line.width = Pt(2.5)
        sr.smooth = True
    if curves_full_data["pred_exact_val"][1] > 1e6:
        xml_node = etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/chart}dispUnits")
        etree.SubElement(xml_node, '{http://schemas.openxmlformats.org/drawingml/2006/chart}builtInUnit').attrib[
            'val'] = 'millions'
        dispUnitsLbl = etree.SubElement(xml_node,
                                        '{http://schemas.openxmlformats.org/drawingml/2006/chart}dispUnitsLbl')
        layout = etree.SubElement(dispUnitsLbl, '{http://schemas.openxmlformats.org/drawingml/2006/chart}layout')

        manual = etree.SubElement(layout, '{http://schemas.openxmlformats.org/drawingml/2006/chart}manualLayout')
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}xMode').attrib['val'] = 'edge'
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}yMode').attrib['val'] = 'edge'
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}x').attrib['val'] = '0.13'
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}y').attrib['val'] = '0.03'

        if model_rep.language == 'russian':
            tx = etree.SubElement(dispUnitsLbl, '{http://schemas.openxmlformats.org/drawingml/2006/chart}tx')
            rich = etree.SubElement(tx, '{http://schemas.openxmlformats.org/drawingml/2006/chart}rich')
            etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')

            p = etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')

            pPr = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')

            r = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
            rpr = etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')

            rpr.attrib['lang'] = 'ru-RU'
            rpr.attrib['dirty'] = '0'

            etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t').text = 'млн'

            etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}endParaRPr').attrib['lang'] = 'en-US'

        chart._chartSpace.xpath("//c:valAx")[0].insert(-1, xml_node)

    # ROI
    chart_data = _chart_data_worker(curves_full_data, f'ROI_{plot_type}', if_exact, n_years)

    x, y = (0.5 - 0.45 + 0.25 + 0.25) * slide_width, slide.shapes[1].top + slide.shapes[1].height
    cx, cy = 0.4 * slide_width, 0.5 * slide_height

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    chart = set_chart_dashes(chart, '_extrapolated')
    chart.series.smooth = True
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.name = model_rep.pptx_cnf['font_family_body']
    chart.legend.font.size = model_rep.pptx_cnf['font_size_small']

    chart = _set_chart_axis(chart, model_rep,
                            if_exact=if_exact,
                            var='roi_axis_y',
                            year=axis_years,
                            plot_type=plot_type,
                            bdg_max=bdg_max,
                            budget_step=budget_step)

    chart = set_chart_colors_line(chart, model_rep, RGBColor(170, 170, 170))

    chart.font.size = model_rep.pptx_cnf['font_size_small']
    chart.font.name = model_rep.pptx_cnf['font_family_body']

    graphic_frame.left = round((0.5 - 0.45 + 0.25 + 0.25) * slide_width)

    for sr in chart.plots[0].series:
        sr.format.line.width = Pt(2.5)
        sr.smooth = True

    return slide


def _chart_data_worker(curves_full_df: pd.DataFrame,
                       curve_type: str,
                       if_exact: bool,
                       n_years: int = 1):
    """
    Prepares already calculated efficiency curves data for specific curve. 

    :param curves_full_df: calculated efficiency curves data
    :param curve_type: 'increment_volume_short', 'increment_value_short', 'profit_short', 'ROI_short'
        , 'increment_volume_long', 'increment_value_long', 'profit_long', 'ROI_long'
    :param if_exact: apply only to the specific time period, without next years
    :param n_years: number of next years accounted in 'long' curves
    :return: ChartData object from python_pptx package
    """

    chart_data = ChartData()
    chart_data.categories = curves_full_df['bdg'].sort_values().unique()

    match curve_type:
        case 'increment_volume_short':
            zero = curves_full_df.loc[curves_full_df['option'] == 'zero', 'pred_exact_vol'].values
            curves_full_df['value'] = curves_full_df['pred_exact_vol'] - zero
        case 'increment_value_short':
            zero = curves_full_df.loc[curves_full_df['option'] == 'zero', 'pred_exact_val'].values
            curves_full_df['value'] = curves_full_df['pred_exact_val'] - zero
        case 'profit_short':
            zero = curves_full_df.loc[curves_full_df['option'] == 'zero', 'pred_exact_val'].values
            curves_full_df['value'] = curves_full_df['pred_exact_val'] - zero - curves_full_df['bdg'] * 1e6
        case 'ROI_short':
            zero = curves_full_df.loc[curves_full_df['option'] == 'zero', 'pred_exact_val'].values
            curves_full_df['value'] = (curves_full_df['pred_exact_val'] - zero) / (curves_full_df['bdg'] * 1e6)
        case 'increment_volume_long':
            zero = curves_full_df.loc[curves_full_df['option'] == 'zero', 'pred_long_vol'].values
            curves_full_df['value'] = curves_full_df['pred_long_vol'] - zero
        case 'increment_value_long':
            zero = curves_full_df.loc[curves_full_df['option'] == 'zero', 'pred_long_val'].values
            curves_full_df['value'] = curves_full_df['pred_long_val'] - zero
        case 'profit_long':
            zero = curves_full_df.loc[curves_full_df['option'] == 'zero', 'pred_long_val'].values
            # fix for correct calcs with if_exact = False
            if if_exact is False:
                curves_full_df['value'] = (curves_full_df['pred_long_val'] - zero) - curves_full_df['bdg'] * 1e6 * n_years
            else:
                curves_full_df['value'] = (curves_full_df['pred_long_val'] - zero) - curves_full_df['bdg'] * 1e6
        case 'ROI_long':
            zero = curves_full_df.loc[curves_full_df['option'] == 'zero', 'pred_long_val'].values
            curves_full_df['value'] = (curves_full_df['pred_long_val'] - zero) / (curves_full_df['bdg'] * 1e6)

    if ("long" in curve_type) & (if_exact is False):
        curves_full_df['value'] = curves_full_df['value'] / n_years

    tmp = curves_full_df.groupby('option').agg({'value': 'nunique'})
    opts = tmp.loc[tmp['value'] > 1].reset_index().option.tolist()
    for opt in opts:
        tmp = curves_full_df.loc[(curves_full_df['option'] == opt)].copy()
        tmp.loc[-1] = 0
        if "ROI" in curve_type:
            tmp.loc[-1] = "#N/A"

        tmp.index = tmp.index + 1
        tmp.sort_index(inplace=True)
        if 'max_obs_budg' in curves_full_df.columns.tolist():
            tmp_obs = tmp.loc[(tmp['bdg'] <= tmp['max_obs_budg'])].copy()
            if len(tmp_obs) > 0:
                chart_data.add_series(opt, tmp_obs['value'])
                last_obs = tmp_obs.iloc[-1]

            tmp_extra = tmp.loc[(tmp['bdg'] > tmp['max_obs_budg'])].copy()
            if len(tmp_extra) > 0:
                if len(tmp_obs) > 0:
                    tmp_extra.loc[-1] = last_obs
                    tmp_extra.index = tmp_extra.index + 1
                    tmp_extra.sort_index(inplace=True)
                chart_data.add_series(opt + '_extrapolated',
                                      pd.concat([pd.Series([''] * (len(tmp_obs) - 1)), tmp_extra['value']], ignore_index=True))

        else:
            chart_data.add_series(opt, tmp['value'])

    return chart_data


def _set_chart_axis(chart,
                    model_rep: "ModelRep",
                    var: str,
                    year: str,
                    budget_step: int | float = 5,
                    bdg_max: int | float = 301,
                    plot_type: str = 'short',
                    if_exact: bool = False):
    """
    Set axis attributes for efficiency curves chart: axis titles, formatting etc.

    :param chart: PPTX chart
    :param model_rep: ModelRep object (reporting settings)
    :param var: kind of plot type; cursed
    :param year: year or years (as string with dash) to be added to titles
    :param budget_step: budget iteration step in millions, defaults to 1 (i.e. 1M)
    :param bdg_max: maximum budget size (all options with larger budgets to be dropped)
    :param plot_type: 'short' or 'long'
    :param if_exact: apply only to the specific time period, without next years
    :return:
    """

    value_axis = chart.value_axis
    value_axis_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == f'efficiency_{plot_type}') & (model_rep.vis_dict['variable'] == var), model_rep.language].iloc[0]
    value_axis_title = value_axis_title.replace('***', year)
    value_axis.axis_title.text_frame.text = value_axis_title
    value_axis.axis_title.text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_small']

    category_axis = chart.category_axis
    category_axis.has_title = True
    if if_exact & (plot_type == 'long'):
        category_axis_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == f'efficiency_{plot_type}') & (
                    model_rep.vis_dict['variable'] == 'axis_x_if_exact_true'), model_rep.language].iloc[0]

        year = year[-4:]
    elif (~if_exact) & (plot_type == 'long'):
        category_axis_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == f'efficiency_{plot_type}') & (
                model_rep.vis_dict['variable'] == 'axis_x_if_exact_false'), model_rep.language].iloc[0]
    else:
        category_axis_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == f'efficiency_{plot_type}') & (model_rep.vis_dict['variable'] == 'axis_x'), model_rep.language].iloc[0]

    category_axis_title = category_axis_title.replace('***', year)
    category_axis.axis_title.text_frame.text = category_axis_title
    category_axis.axis_title.text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_small']

    etree.SubElement(chart._chartSpace.xpath("//c:valAx")[0], '{http://schemas.openxmlformats.org/drawingml/2006/chart}crossBetween').attrib['val'] = 'midCat'

    skip = str(round((bdg_max / budget_step) / 5))

    if int(skip) < 1:
        skip = "1"

    etree.SubElement(chart._chartSpace.xpath("//c:catAx")[0], '{http://schemas.openxmlformats.org/drawingml/2006/chart}tickLblSkip').attrib['val'] = skip
    etree.SubElement(chart._chartSpace.xpath("//c:catAx")[0], '{http://schemas.openxmlformats.org/drawingml/2006/chart}tickMarkSkip').attrib['val'] = skip

    return chart


