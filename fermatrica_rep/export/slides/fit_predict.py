"""
Generate slide with fit and prediction plot. Could be used both in retro analysis and in option reporting,
for superbrand as a whole and for specific number of variables known as "SKU"
"""


import pandas as pd
import numpy as np
import copy

from pptx.presentation import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import lxml.etree as etree

from fermatrica_utils import date_to_period
from fermatrica.model.model import Model

from fermatrica_rep.export.basics import set_table_border, fill_table
from fermatrica_rep.meta_model.model_rep import ModelRep
from fermatrica_rep.options.define import OptionSettings


def create(prs: Presentation,
           model: "Model",
           model_rep: "ModelRep",
           dt_pred: pd.DataFrame,
           opt_set: OptionSettings,
           option_name: str,
           option: dict | None = None,
           period: str = 'day',  # 'day', 'week', 'month', 'quarter', 'year'
           group_var: list | tuple = ('superbrand',),
           plot_type: str = 'brand',
           bs_key_filter: list | tuple | None = None,
           show_future: bool = False
           ):
    """
    Create fit-predict slide and add to `prs` PPTX presentation. Fit-predict slide contains multiple objects,
    so number of worker functions are called from this function.

    Main objects of the slide (aside title etc.):

    1. Volume fit-predict plot
    2. Value fit-predict plot
    3. Volume yearly table
    4. Value yearly table
    5. Option table

    :param prs: Presentation object from python_pptx package
    :param model: Model object
    :param model_rep: ModelRep object (reporting settings)
    :param dt_pred: prediction data
    :param opt_set: OptionSetting object containing calculate settings
    :param option_name: name of the option to be reported. To be used only as title here, no impact on program behaviour
    :param option: dictionary containing option to calculate or None for retro-analysis
    :param period: time period / interval to group by: 'day', 'week', 'month', 'quarter', 'year'
    :param group_var: group entities by variables (list or tuple of strings)
    :param plot_type: 'retro', 'brand' or 'sku'. Used for naming only
    :param bs_key_filter: list or tuple of 'bs_key' values to preserve
    :param show_future: show future periods or not
    :return: Presentation object from python_pptx package
    """

    language = model_rep.language
    obs_name = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "predict_table") & (
                model_rep.vis_dict['variable'] == "observed"), language].iloc[0]
    pred_name = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "predict_table") & (
                model_rep.vis_dict['variable'] == "predicted"), language].iloc[0]

    superbrand = model.conf.target_superbrand

    if hasattr(model.conf, 'display_units'):
        units = model.conf.display_units
    else:
        units = 'packs'

    axis_title_vol = model_rep.vis_dict.loc[
        (model_rep.vis_dict['section'] == "shared") & (model_rep.vis_dict['variable'] == units), language].iloc[0]
    axis_title_val = model_rep.vis_dict.loc[
        (model_rep.vis_dict['section'] == "shared") & (model_rep.vis_dict['variable'] == "rub_short"), language].iloc[0]
    chart_title_vol = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "plot_predict") & (
                model_rep.vis_dict['variable'] == "volume_title_total"), language].iloc[0]
    chart_title_vol = chart_title_vol.replace('***', model_rep.vis_dict.loc[
        (model_rep.vis_dict['section'] == "shared") & (model_rep.vis_dict['variable'] == units), language].iloc[0])

    chart_title_val = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "plot_predict") & (
                model_rep.vis_dict['variable'] == "value_title"), language].iloc[0]

    # !!! temporary fix
    dt_pred = dt_pred[dt_pred['predicted'] >= 0]
    # dt_pred[dt_pred['predicted'] < 0]['predicted'] = pd.NA

    # load data
    ds_melt_vol, ds_melt_val = _line_chart_data(model=model
                                                , dt_pred=dt_pred
                                                , model_rep=model_rep
                                                , period=period
                                                , show_future=show_future
                                                , group_var=group_var
                                                , bs_key_filter=bs_key_filter)

    # prepare slide in presentation

    slide_layout = model_rep.pptx_cnf['Blank_slide']

    slide_height = model_rep.pptx_cnf['slide_height']
    slide_width = model_rep.pptx_cnf['slide_width']

    slide = prs.slides.add_slide(slide_layout)

    # Volume

    vol_chart = _fit_chart_create(slide,
                                  model_rep=model_rep,
                                  model=model,
                                  ds_melt_prep=ds_melt_vol,
                                  axis_title=axis_title_vol,
                                  chart_title=chart_title_vol,
                                  x_offset=0.45,
                                  plot_color=(255, 0, 0))

    # Value

    val_chart = _fit_chart_create(slide,
                                  model_rep=model_rep,
                                  model=model,
                                  ds_melt_prep=ds_melt_val,
                                  axis_title=axis_title_val,
                                  chart_title=chart_title_val,
                                  x_offset=0,
                                  plot_color=(0, 0, 255))

    # Tables

    dt_pred_prep = \
    dt_pred.assign(year=dt_pred['date'].dt.year, observed_val=dt_pred['observed'] * dt_pred[model.conf.price_var]
                   , predicted_val=dt_pred['predicted'] * dt_pred[model.conf.price_var]). \
        loc[(dt_pred['bs_key'].isin(bs_key_filter)) & (dt_pred['listed'].isin([2, 3, 4]))]. \
        groupby('year')[['observed', 'predicted', 'observed_val', 'predicted_val']].sum().reset_index()
    dt_pred_prep['err'] = dt_pred_prep['predicted'] / dt_pred_prep['observed'] - 1
    dt_pred_prep['err'] = [f"{x:.1%}" if abs(x) != np.inf else '-' for x in dt_pred_prep['err']]
    dt_pred_prep['err_val'] = dt_pred_prep['predicted_val'] / dt_pred_prep['observed_val'] - 1
    dt_pred_prep['err_val'] = [f"{x:.1%}" if abs(x) != np.inf else '-' for x in dt_pred_prep['err_val']]

    for col in dt_pred_prep[['observed', 'predicted', 'observed_val', 'predicted_val']]:
        dt_pred_prep[col] = ['{:,.0f}'.format(x).replace(',', ' ') if x >= 0 else '0' for x in dt_pred_prep[col]]

    dt_pred_prep_vol = dt_pred_prep[['year', 'observed', 'predicted', 'err']]. \
        rename(columns={'year': model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'year', language].iloc[0],
                        'observed':
                            model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'observed', language].iloc[0],
                        'predicted':
                            model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'predicted', language].iloc[0],
                        'err': model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'err', language].iloc[0]
                        })

    dt_pred_prep_val = dt_pred_prep[['year', 'observed_val', 'predicted_val', 'err_val']]. \
        rename(columns={'year': model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'year', language].iloc[0],
                        'observed_val':
                            model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'observed', language].iloc[0],
                        'predicted_val':
                            model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'predicted', language].iloc[0],
                        'err_val': model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'err', language].iloc[0]
                        })

    rows_pred = len(dt_pred_prep.index) + 1
    cell_width = round(1/12 * slide_width)
    cell_height = round(1/40 * slide_height)

    # Volume table

    x, y = 0.5 * slide_width - 0.45 * slide_width, slide.shapes[0].top + slide.shapes[0].height + 0.5 * slide_height
    cx, cy = 4 * cell_width, rows_pred * cell_height

    x, y, cx, cy = int(round(x)), int(round(y)), int(round(cx)), int(round(cy))

    shape = slide.shapes.add_table(rows_pred, 4, x, y, cx, cy)
    table = shape.table
    fill_table(table, dt_pred_prep_vol)

    for ind, cell in enumerate(table.iter_cells()):
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = model_rep.pptx_cnf['font_size_small']
            paragraph.font.name = model_rep.pptx_cnf['font_family_body']
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    table.columns[0].width = round(table.columns[0].width * .75)
    table.columns[1].width = round(table.columns[1].width * 1.25)
    table.columns[2].width = round(table.columns[2].width * 1.25)
    table.columns[3].width = round(table.columns[3].width * .75)

    set_table_border(table,
                     border_color='595959',
                     border_width='25000')

    # Value table

    x, y = 0.5 * slide_width + 0.45 * slide_width - 4 * cell_width, slide.shapes[0].top + slide.shapes[0].height + 0.5 * slide_height
    cx, cy = 4 * cell_width, rows_pred * cell_height

    x, y, cx, cy = int(round(x)), int(round(y)), int(round(cx)), int(round(cy))

    graphic_frame = slide.shapes.add_table(rows_pred, 4, x, y, cx, cy)
    table = graphic_frame.table

    fill_table(table, dt_pred_prep_val)

    for ind, cell in enumerate(table.iter_cells()):
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = model_rep.pptx_cnf['font_size_small']
            paragraph.font.name = model_rep.pptx_cnf['font_family_body']
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    table.columns[0].width = round(table.columns[0].width * .75)
    table.columns[1].width = round(table.columns[1].width * 1.25)
    table.columns[2].width = round(table.columns[2].width * 1.25)
    table.columns[3].width = round(table.columns[3].width * .75)

    table_width = 0
    for col in table.columns:
        table_width += col.width

    table_height = 0
    for row in table.rows:
        table_height += row.height

    graphic_frame.left = round(0.5 * slide_width + 0.45 * slide_width - table_width)

    set_table_border(table,
                     border_color='595959',
                     border_width='25000')

    # Option table

    if option is not None:
        slide = _option_table_create(model_rep=model_rep
                                     , slide=slide
                                     , option=option
                                     , cell_width=cell_width
                                     , cell_height=cell_height)

    # Title

    if opt_set.date_start.year in dt_pred_prep_vol[
        model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'year', language].iloc[0]].values:
        vol_kpi = dt_pred_prep_vol[dt_pred_prep_vol[model_rep.vis_dict.loc[
            model_rep.vis_dict['variable'] == 'year', language].iloc[0]] == opt_set.date_start.year][
            model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'predicted', language].iloc[0]].iloc[0]
        vol_kpi = _num_kpi_prep(vol_kpi, model_rep)
    else:
        vol_kpi = '0'

    if opt_set.date_start.year in dt_pred_prep_val[
        model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'year', language].iloc[0]].values:
        val_kpi = dt_pred_prep_val[dt_pred_prep_val[model_rep.vis_dict.loc[
            model_rep.vis_dict['variable'] == 'year', language].iloc[0]] == opt_set.date_start.year][
            model_rep.vis_dict.loc[model_rep.vis_dict['variable'] == 'predicted', language].iloc[0]].iloc[0]
        val_kpi = _num_kpi_prep(val_kpi, model_rep)
    else:
        val_kpi = '0'

    if plot_type in ['brand', 'sku']:
        slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'option.detail.sales') &
                                             (model_rep.vis_dict['variable'] == 'title'), language].iloc[0]
    else:
        slide_title = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'predict_table') &
                                             (model_rep.vis_dict['variable'] == 'title'), language].iloc[0]

    slide_title = slide_title.replace('***id', option_name)
    slide_title = slide_title.replace('***year', str(opt_set.date_start.year))
    slide_title = slide_title.replace('***volume_sales', vol_kpi)
    slide_title = slide_title.replace('***value_sales', val_kpi)
    slide_title = slide_title.replace('***units', model_rep.vis_dict.loc[
        (model_rep.vis_dict['section'] == "shared") & (model_rep.vis_dict['variable'] == units), language].iloc[0])

    slide.shapes[0].text_frame.text = slide_title
    slide.shapes[0].text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_gross']
    slide.shapes[0].text_frame.paragraphs[0].font.name = model_rep.pptx_cnf['font_family_header']

    return prs


"""
Visuals
"""


def _fit_chart_create(slide,
                      model: "Model",
                      model_rep: "ModelRep",
                      ds_melt_prep: pd.DataFrame,
                      axis_title: str,
                      chart_title: str,
                      x_offset: float,
                      plot_color: tuple = (0, 0, 0)):
    """
    Creates fit-predict plot for `fermatrica_rep.export.slides.fit_predict.create()` slide.
    Call it twice with specific settings to make volume and value plots respectively.

    :param slide: Slide object from python_pptx package
    :param model: Model object
    :param model_rep: ModelRep object (reporting settings)
    :param ds_melt_prep: prepared data
    :param axis_title: Y title as string
    :param chart_title: plot title as string
    :param x_offset: place plot horizontally (looks like it is measured from the right?)
    :param plot_color: RGB color as tuple
    :return: Chart object from python_pptx package
    """

    slide_width = model_rep.pptx_cnf['slide_width']
    slide_height = model_rep.pptx_cnf['slide_height']

    superbrand = model.conf.target_superbrand

    language = model_rep.language
    obs_name = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "predict_table") &
                                      (model_rep.vis_dict['variable'] == "observed"), language].iloc[0]
    pred_name = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "predict_table") &
                                       (model_rep.vis_dict['variable'] == "predicted"), language].iloc[0]

    chart_data_val = ChartData()
    chart_data_val.categories = ds_melt_prep[ds_melt_prep['variable'] == obs_name]["date"]
    chart_data_val.add_series(name=obs_name,
                              values=ds_melt_prep.loc[ds_melt_prep["variable"] == obs_name, "value"])
    chart_data_val.add_series(name=pred_name,
                              values=ds_melt_prep.loc[ds_melt_prep["variable"] == pred_name, "value"])

    x, y = 0.5 * slide_width - x_offset * slide_width, slide.shapes[0].top + slide.shapes[0].height
    cx, cy = 0.45 * slide_width, 0.5 * slide_height

    x, y, cx, cy = int(round(x)), int(round(y)), int(round(cx)), int(round(cy))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_val
    ).chart

    chart.plots[0].series[0].format.line.color.rgb = RGBColor(0, 0, 0)
    chart.plots[0].series[1].format.line.color.rgb = RGBColor(*plot_color)

    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.chart_title.text_frame.text = f"{superbrand.capitalize()}: {chart_title}"
    chart.chart_title.text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_main']

    category_axis = chart.category_axis
    category_axis.has_minor_gridlines = True
    category_axis.has_major_gridlines = False

    value_axis = chart.value_axis
    value_axis.axis_title.text_frame.text = axis_title
    value_axis.axis_title.text_frame.paragraphs[0].font.size = model_rep.pptx_cnf['font_size_main']

    value_axis.minimum_scale = 0
    value_axis.maximum_scale = round(1.2 * max(ds_melt_prep["value"]))

    value_axis.has_minor_gridlines = True
    value_axis.has_major_gridlines = False

    chart.font.size = model_rep.pptx_cnf['font_size_small']
    chart.font.name = model_rep.pptx_cnf['font_family_body']

    chart.has_legend = True
    chart.legend.include_in_layout = False

    if chart.value_axis.maximum_scale > 1000000:

        xml_node = etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/chart}dispUnits")
        etree.SubElement(xml_node, '{http://schemas.openxmlformats.org/drawingml/2006/chart}builtInUnit').attrib[
            'val'] = 'millions'
        dispUnitsLbl = etree.SubElement(xml_node,
                                        '{http://schemas.openxmlformats.org/drawingml/2006/chart}dispUnitsLbl')
        layout = etree.SubElement(dispUnitsLbl, '{http://schemas.openxmlformats.org/drawingml/2006/chart}layout')

        manual = etree.SubElement(layout, '{http://schemas.openxmlformats.org/drawingml/2006/chart}manualLayout')
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}xMode').attrib['val'] = 'edge'
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}yMode').attrib['val'] = 'edge'
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}x').attrib['val'] = '0.13'
        etree.SubElement(manual, '{http://schemas.openxmlformats.org/drawingml/2006/chart}y').attrib['val'] = '0.13'

        if language == 'russian':
            tx = etree.SubElement(dispUnitsLbl, '{http://schemas.openxmlformats.org/drawingml/2006/chart}tx')
            rich = etree.SubElement(tx, '{http://schemas.openxmlformats.org/drawingml/2006/chart}rich')
            etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')

            p = etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')

            pPr = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')

            r = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
            rpr = etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')

            rpr.attrib['lang'] = 'ru-RU'
            rpr.attrib['dirty'] = '0'

            etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t').text = 'млн'

            etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}endParaRPr').attrib['lang'] = 'en-US'

        chart._chartSpace.xpath("//c:valAx")[0].insert(-1, xml_node)

    return chart


def _option_table_create(slide
                         , model_rep: "ModelRep"
                         , option: dict
                         , cell_width: int | float
                         , cell_height: int | float):
    """
    Creates table with option / scenario / budget specifications. To be placed in the middle of the
    bottom of the slide.

    :param slide: Slide object from python_pptx package
    :param model_rep: ModelRep object (reporting settings)
    :param option: dictionary containing option to calculate or None for retro-analysis
    :param cell_width: width of the cell in the table
    :param cell_height: height of the cell in the table
    :return: Slide object from python_pptx package
    """

    language = model_rep.language

    slide_height = model_rep.pptx_cnf['slide_height']
    slide_width = model_rep.pptx_cnf['slide_width']

    #

    bdg_col = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'predict_costs_table') &
                                     (model_rep.vis_dict['variable'] == 'value'), language].iloc[0]
    invest_col = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'predict_costs_table') &
                                        (model_rep.vis_dict['variable'] == 'label'), language].iloc[0].replace(',', '')

    option_df = pd.DataFrame.from_dict(option.items()).rename(columns={0: invest_col,
                                                                       1: bdg_col})

    option_df = option_df[option_df[invest_col] != "text_flags"].reset_index().drop('index', axis=1)
    option_df = option_df[pd.to_numeric(option_df[bdg_col], errors='coerce').notnull()]
    option_df[bdg_col] = [f"{x[1][1]:.0%}" if 'bdg' not in x[1][0] else str(x[1][1]) for x in option_df.iterrows()]

    option_df['tmp'] = option_df[invest_col].str.replace('bdg_', '')
    option_df = option_df.merge(
        model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == 'costs_type'), ['variable', language]]
        , how='left', left_on='tmp', right_on='variable')

    option_df = option_df.groupby([invest_col, bdg_col]).first().reset_index()
    option_df.loc[option_df[language].notna(), invest_col] = option_df.loc[option_df[language].notna(), language]

    option_df.drop(language, axis=1, inplace=True)
    option_df.drop('tmp', axis=1, inplace=True)
    option_df.drop('variable', axis=1, inplace=True)

    rows_opt = len(option_df) + 1

    x, y = 0.5 * slide_width - cell_width, slide.shapes[0].top + slide.shapes[0].height + 0.5 * slide_height
    cx, cy = 2 * cell_width, rows_opt * cell_height

    x, y, cx, cy = int(round(x)), int(round(y)), int(round(cx)), int(round(cy))

    graphic_frame = slide.shapes.add_table(rows_opt, 2, x, y, cx, cy)
    table = graphic_frame.table

    #

    fill_table(table, option_df)

    for ind, cell in enumerate(table.iter_cells()):
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = model_rep.pptx_cnf['font_size_small']
            paragraph.font.name = model_rep.pptx_cnf['font_family_body']
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    table.columns[0].width = round(table.columns[0].width * 1.25)
    table.columns[1].width = round(table.columns[1].width * 1.25)

    graphic_frame.left = round(0.5 * slide_width - table.columns[0].width)

    return slide


def _num_kpi_prep(string: str
                  , model_rep: "ModelRep"):
    """
    Format KPI value to be placed into slide title of fit-predict slide.

    :param string: numeric passes as string (extracted from data prepared for table)
    :param model_rep: ModelRep object (reporting settings)
    :return:
    """

    num = int(string.replace(' ', ''))
    match num:
        case num if 1e+3 <= num < 1e+6:
            return str(round(num / 1e+3)) + ' ' + model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "shared") &
                                                                         (model_rep.vis_dict['variable'] == "thousand_short"),
                                                                         model_rep.language].iloc[0]
        case num if 1e+6 <= num < 1e+9:
            return str(round(num / 1e+6, 2)).replace('.', ',') + ' ' + model_rep.vis_dict.loc[
                (model_rep.vis_dict['section'] == "shared") & (model_rep.vis_dict['variable'] == "million_short"),
                model_rep.language].iloc[0] if model_rep.language == 'russian' else \
                str(round(num / 1e+6, 2)) + ' ' + model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "shared") &
                                                                         (model_rep.vis_dict['variable'] == "million_short"),
                                                                         model_rep.language].iloc[0]

        case num if num >= 1e+9:
            return str(round(num / 1e+9, 2)).replace('.', ',') + ' ' + model_rep.vis_dict.loc[
                (model_rep.vis_dict['section'] == "shared") & (model_rep.vis_dict['variable'] == "billion_short"),
                model_rep.language].iloc[0] if model_rep.language == 'russian' else \
                str(round(num / 1e+9, 2)) + ' ' + model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "shared") &
                                                                         (model_rep.vis_dict['variable'] == "billion_short"),
                                                                         model_rep.language].iloc[0]
        case _:
            return '0'


"""
Data
"""


def _line_chart_data(model: "Model",
                     dt_pred: pd.DataFrame,
                     model_rep: "ModelRep",
                     period: str = 'day',  # 'day', 'week', 'month', 'quarter', 'year'
                     group_var: list | tuple = ('superbrand', ),
                     bs_key_filter: list | tuple | None = None,
                     show_future: bool = False,
                     ):
    """
    Convert 'common' prediction dataset to dataset prepared to be plotted as fit-predict.

    :param model: Model object
    :param dt_pred: prediction data
    :param model_rep: ModelRep object (reporting settings)
    :param period: time period / interval to group by: 'day', 'week', 'month', 'quarter', 'year'
    :param group_var: group entities by variables (list or tuple of strings)
    :param bs_key_filter: list or tuple of 'bs_key' values to preserve
    :param show_future: show future periods or not
    :return:
    """

    language = model_rep.language
    obs_name = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "predict_table") & (
                model_rep.vis_dict['variable'] == "observed"), language].iloc[0]
    pred_name = model_rep.vis_dict.loc[(model_rep.vis_dict['section'] == "predict_table") & (
                model_rep.vis_dict['variable'] == "predicted"), language].iloc[0]

    # load data
    ds, ds_melt = _fit_main_data_melt(model=model
                                      , dt_pred=dt_pred
                                      , period=period
                                      , group_var=group_var
                                      , bs_key_filter=bs_key_filter
                                      , show_future=show_future)

    # Volume
    ds_melt_vol = ds_melt[(ds_melt['variable'] == 'observed') | (ds_melt['variable'] == 'predicted')]

    ds_melt_vol.loc[ds_melt_vol['variable'] == 'observed', 'variable'] = obs_name
    ds_melt_vol.loc[ds_melt_vol['variable'] == 'predicted', 'variable'] = pred_name

    ds_melt_vol.loc[:, 'group_id'] = ds_melt_vol.loc[:, 'group_id'].str.title()

    # Value
    ds_melt_val = ds_melt[(ds_melt['variable'] == 'observed_value') | (ds_melt['variable'] == 'predicted_value')]

    ds_melt_val.loc[ds_melt_val['variable'] == 'observed_value', 'variable'] = obs_name
    ds_melt_val.loc[ds_melt_val['variable'] == 'predicted_value', 'variable'] = pred_name

    ds_melt_val.loc[:, 'group_id'] = ds_melt_val.loc[:, 'group_id'].str.title()

    return ds_melt_vol, ds_melt_val


def _fit_main_data_melt(model: "Model"
                        , dt_pred: pd.DataFrame
                        , period: str = 'day'  # 'day', 'week', 'month', 'quarter', 'year'
                        , group_var: list | tuple = ('superbrand', )
                        , bs_key_filter: list | tuple | None = None
                        , show_future: bool = False
                        ) -> (pd.DataFrame, pd.DataFrame):
    """
    Melt prediction data to make it suitable to PPTX charts in fit-predict slide.

    :param model: Model object
    :param dt_pred: prediction data
    :param period: time period / interval to group by: 'day', 'week', 'month', 'quarter', 'year'
    :param group_var: group entities by variables (list or tuple of strings)
    :param bs_key_filter: list or tuple of 'bs_key' values to preserve
    :param show_future: show future periods or not
    :return:
    """

    # filter

    if show_future:
        ds = dt_pred[dt_pred['listed'].isin([2, 3, 4])]
    else:
        ds = dt_pred[dt_pred['listed'].isin([2, 3])]

    if bs_key_filter is not None:
        ds = ds[ds['bs_key'].isin(bs_key_filter)]

    ds = copy.deepcopy(ds)

    # get values

    ds['observed_value'] = ds['observed'] * ds[model.conf.price_var]
    ds['predicted_value'] = ds['predicted'] * ds[model.conf.price_var]

    if hasattr(model.conf, 'conversion_var') and model.conf.conversion_var is not None:
        ds['observed_value'] = ds['observed_value'] * ds[model.conf.conversion_var]
        ds['predicted_value'] = ds['predicted_value'] * ds[model.conf.conversion_var]

    # reduce

    ds['date'] = date_to_period(ds['date'], period)

    if type(group_var) == tuple:
        group_var = list(group_var)

    group_var_ext = copy.deepcopy(group_var)

    group_var_ext.extend(['date', 'listed'])

    cols = copy.deepcopy(group_var_ext)
    cols.extend(['observed', 'predicted', 'observed_value', 'predicted_value'])

    ds = ds[cols].groupby(group_var_ext).sum().reset_index()

    # melt

    ds_melt = pd.melt(ds, id_vars=group_var_ext)

    ds['group_id'] = ''
    ds_melt['group_id'] = ''

    for i in group_var:
        ds['group_id'] = ds['group_id'] + ' ' + ds[i].astype('str')
        ds_melt['group_id'] = ds_melt['group_id'] + ' ' + ds_melt[i].astype('str')

    return ds, ds_melt


