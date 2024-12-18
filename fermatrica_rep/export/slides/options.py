"""
Generate slide(s) with comparative table containing summaries by number of options,
defined in `options_m` dictionary.
"""


import copy
import pandas as pd
import numpy as np

from pptx.presentation import Presentation
from pptx.dml.color import RGBColor

from fermatrica_rep.export.basics import set_table_border, adjust_table_width, table_text_format
from fermatrica_rep.options.define import OptionSettings
import fermatrica_rep.options.calc as calc


def create(prs: Presentation,
           model: "Model | list",
           model_rep: "ModelRep | list",
           dt_pred: pd.DataFrame | list,
           ds: pd.DataFrame | list,
           options_m: dict,
           opt_set: "OptionSettings",
           bs_key_filter: list | tuple | None = None,
           if_exact: bool = True
           ):
    """
    Generate slide(s) with comparative table containing summaries by number of options,
    defined in `options_m` dictionary. If number of options in `options_m` is more than 4 or 5
    or something, more than 1 slide to be generated to make table neat, readable and user-friendly.

    There are two tables to be placed onto slide(s) one below another:

    1. Option (budget) table
    2. Summary table

    :param prs: Presentation object from python_pptx package
    :param model: Model object or list of Model objects
    :param model_rep: ModelRep object (reporting settings) or list ModelRep objects
    :param dt_pred: prediction data or list of prediction datas
    :param ds: main dataset or list of main datasets
    :param options_m: dictionary of dictionaries defining options to calculate
    :param opt_set: OptionSetting object containing calculate settings
    :param bs_key_filter: list or tuple of 'bs_key' values to preserve
    :param if_exact: apply only to the specific time period, without next years
    :return: Presentation object from python_pptx package
    """

    if isinstance(model_rep, list):
        model_rep_main = model_rep[-1]
        dt_pred = dt_pred[-1]
        price_var = model[-1].conf.price_var
    else:
        model_rep_main = model_rep
        price_var = model.conf.price_var

    language = model_rep_main.language
    slide_width = model_rep_main.pptx_cnf['slide_width']
    slide_height = model_rep_main.pptx_cnf['slide_height']

    # slide_layout = [x for x in prs.slide_masters[1].slide_layouts if x.name == "Blank_slide"][0]
    slide_layout = model_rep_main.pptx_cnf['Blank_slide']

    dt_pred_prep = dt_pred.assign(year=dt_pred['date'].dt.year,
                                  observed_val=dt_pred['observed'] * dt_pred[price_var],
                                  predicted_val=dt_pred['predicted'] * dt_pred[price_var]). \
        loc[(dt_pred['bs_key'].isin(bs_key_filter)) & (dt_pred['listed'].isin([2, 3, 4]))]. \
        groupby('year')[['observed', 'predicted', 'observed_val', 'predicted_val']].sum().reset_index()

    pred_year_start = dt_pred_prep[dt_pred_prep['observed'] == 0]['year'].min()
    pred_year_end = dt_pred_prep['year'].max()

    options_names = copy.deepcopy(options_m)
    options_names = list(options_names)[1:]
    options_names = [options_names[n:n+3] for n in range(0, len(options_names), 3)]

    # temporary
    if language == 'english':
        opt_str = 'Option'
    else:
        opt_str = 'Опция'

    for options_names_subset in options_names:
        slide = prs.slides.add_slide(slide_layout)

        slide_title = model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'option_compare') & (
                model_rep_main.vis_dict['variable'] == 'title'), language].iloc[0]
        slide.shapes[0].text_frame.text = slide_title
        slide.shapes[0].text_frame.paragraphs[0].font.size = model_rep_main.pptx_cnf['font_size_gross']
        slide.shapes[0].text_frame.paragraphs[0].font.name = model_rep_main.pptx_cnf['font_family_header']

        # options text boxes
        x, y = (0.05 + 0.9 * 2/10) * slide_width, slide.shapes[0].top + slide.shapes[0].height + 0.05 * slide_height
        cx, cy = 0.9 * 2/10 * slide_width, 0.05 * slide_height

        x, y, cx, cy = int(round(x)), int(round(y)), int(round(cx)), int(round(cy))

        textbox = slide.shapes.add_textbox(x, y, cx, cy)
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.text = f"{opt_str} 1. {list(options_m)[0]}"
        paragraph.font.bold = True
        paragraph.font.size = model_rep_main.pptx_cnf["font_size_main"]
        paragraph.font.name = model_rep_main.pptx_cnf["font_family_body"]
        textbox.text_frame.word_wrap = True

        for n, option in enumerate(options_names_subset):
            x, y = (0.05 + 0.9 * 2 * (n+2) / 10) * slide_width, slide.shapes[0].top + slide.shapes[0].height + 0.05 * slide_height
            cx, cy = 0.9 * 2/10 * slide_width, 0.05 * slide_height
            textbox = slide.shapes.add_textbox(x, y, cx, cy)
            paragraph = textbox.text_frame.paragraphs[0]
            paragraph.text = f"{opt_str} {list(options_m).index(option) + 1}. {option}"
            paragraph.font.bold = True
            paragraph.font.size = model_rep_main.pptx_cnf["font_size_main"]
            paragraph.font.name = model_rep_main.pptx_cnf["font_family_body"]
            textbox.text_frame.word_wrap = True

        # prepare data
        opt_table_data, calc_table_data = _tables_data_worker(model=model,
                                                              model_rep=model_rep,
                                                              ds=ds,
                                                              opt_set=opt_set,
                                                              options_names_subset=options_names_subset,
                                                              options_m=options_m,
                                                              pred_year_start=pred_year_start,
                                                              pred_year_end=pred_year_end,
                                                              if_exact=if_exact,
                                                              hide_rows=['text_flags', 'google', 'price'])

        # ----- option table -----

        # construct table

        n_rows = len(opt_table_data.iloc[:, 0]) + 1
        n_cols = 3 + 2 * len(options_names_subset)

        cell_width = round(1/10 * slide_width)
        cell_height = round(1/28 * slide_height)

        x, y = (0.5 - 0.45) * slide_width, slide.shapes[1].top + slide.shapes[1].height
        cx, cy = n_cols * cell_width, n_rows * cell_height

        x, y, cx, cy = int(round(x)), int(round(y)), int(round(cx)), int(round(cy))

        opt_shape = slide.shapes.add_table(n_rows, n_cols, x, y, cx, cy)
        opt_table = opt_shape.table

        # set borders

        set_table_border(opt_table,
                         border_color='595959',
                         border_width='25000')

        # fill table with data

        opt_table = _fill_table(opt_table,
                                ds=opt_table_data,
                                pred_year_start=pred_year_start,
                                pred_year_end=pred_year_end,
                                last_row_bold=True)

        # set font, size, position

        opt_table = table_text_format(opt_table, model_rep_main)

        # adjust table size
        opt_table = adjust_table_width(table=opt_table,
                                       slide_width=slide_width)

        # ----- calc table -----

        # prepare data
        # construct table

        n_rows = 7
        n_cols = 3 + 2 * len(options_names_subset)

        x, y = (0.5 - 0.45) * slide_width, opt_shape.top + opt_shape.height + 0.05 * slide_height
        cx, cy = n_cols * cell_width, n_rows * cell_height

        x, y, cx, cy = int(round(x)), int(round(y)), int(round(cx)), int(round(cy))

        calc_shape = slide.shapes.add_table(n_rows, n_cols, x, y, cx, cy)
        calc_table = calc_shape.table

        # set borders

        set_table_border(calc_table,
                         border_color='595959',
                         border_width='25000')

        # fill table with data

        calc_table = _fill_table(calc_table,
                                 ds=calc_table_data,
                                 pred_year_start=pred_year_start,
                                 pred_year_end=pred_year_end)

        # set font, size, position
        calc_table = table_text_format(calc_table, model_rep_main)

        # adjust table size
        calc_table = adjust_table_width(table=calc_table,
                                        slide_width=slide_width)

    return prs


def _tables_data_worker(model: "Model | list",
                        model_rep: "ModelRep | list",
                        ds: pd.DataFrame | list,
                        opt_set: "OptionSettings",
                        options_m: dict,
                        options_names_subset: list,
                        pred_year_start: int,
                        pred_year_end: int,
                        if_exact: bool = True,
                        hide_rows: list | tuple = ('text_flags',)):
    """
    Calculate subset of options and fill the table with option data and summaries for one slide only.
    I.e. all options to be calculated are split into batches, so every batch corresponds one slide.
    More options -> more batches -> more slides.

    First option (zero) is duplicated on every slide to facilitate option compare by the end user.

    :param model: Model object or list of Model objects
    :param model_rep: ModelRep object (reporting settings) or list ModelRep objects
    :param ds: main dataset or list of main datasets
    :param opt_set: OptionSetting object containing calculate settings
    :param options_m: dictionary of dictionaries defining options to calculate
    :param options_names_subset:
    :param pred_year_start:
    :param pred_year_end:
    :param if_exact:
    :param hide_rows:
    :return:
    """

    if isinstance(model_rep, list):
        model_rep_main = model_rep[-1]
    else:
        model_rep_main = model_rep

    # ---------- options table ----------

    language = model_rep_main.language

    invest_col = model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'option_compare') &
                                        (model_rep_main.vis_dict['variable'] == 'costs_title'), model_rep_main.language].iloc[0]

    options_df = pd.DataFrame.from_dict(options_m['zero'].items()).rename(columns={0: invest_col,
                                                                                   1: 'zero_short'})

    options_df = options_df[~options_df[invest_col].isin(hide_rows)].reset_index().drop('index', axis=1)
    options_df = options_df[pd.to_numeric(options_df['zero_short'], errors='coerce').notnull()]
    options_df['zero_long'] = [f"{x[1][1]:.0%}" if 'bdg' not in x[1][0]
                               else str((pred_year_end - pred_year_start + 1) * x[1][1])
                               for x in options_df.iterrows()]
    options_df['zero_short'] = [f"{x[1][1]:.0%}" if 'bdg' not in x[1][0] else str(x[1][1])
                                for x in options_df.iterrows()]

    for option_name in options_names_subset:
        tmp_df = pd.DataFrame.from_dict(options_m[option_name].items()).rename(columns={0: 'tmp',
                                                                                        1: f'{option_name}_short'})
        tmp_df = tmp_df[pd.to_numeric(tmp_df[f'{option_name}_short'], errors='coerce').notnull()]
        tmp_df[f'{option_name}_long'] = [f"{x[1][1]:.0%}" if 'bdg' not in x[1][0]
                                         else str((pred_year_end - pred_year_start + 1) * x[1][1])
                                         for x in tmp_df.iterrows()]
        tmp_df[f'{option_name}_short'] = [f"{x[1][1]:.0%}" if 'bdg' not in x[1][0] else str(x[1][1])
                                          for x in tmp_df.iterrows()]

        options_df = options_df.merge(tmp_df, how='left', left_on=invest_col, right_on='tmp')
        options_df.drop('tmp', axis=1, inplace=True)

    group_cols = list(options_df)

    options_df['tmp'] = options_df[invest_col].str.replace('bdg_', '')
    options_df = options_df.merge(
        model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'costs_type'), ['variable', language]]
        , how='left', left_on='tmp', right_on='variable')

    options_df = options_df.groupby(group_cols).first().reset_index()
    options_df.loc[options_df[language].notna(), invest_col] = options_df.loc[options_df[language].notna(), language]

    options_df.drop(language, axis=1, inplace=True)
    options_df.drop('variable', axis=1, inplace=True)
    options_df.drop('tmp', axis=1, inplace=True)

    # drop rows with all zeros
    options_df = options_df.loc[(options_df[group_cols[1:]] != '0').any(axis=1)]

    # tmp fix to display total bdg in last row
    row_to_move = options_df.iloc[0]
    options_df = options_df.drop(0)
    options_df = pd.concat([options_df, row_to_move.to_frame().T])

    options_df = options_df.reset_index().drop('index', axis=1)

    # ---------- calc table ----------

    sales_col_name = model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'option_compare') &
                                            (model_rep_main.vis_dict['variable'] == 'sales_title'), model_rep_main.language].iloc[0]
    sales_vol = model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'option_compare_sales') &
                                       (model_rep_main.vis_dict['variable'] == 'res'), model_rep_main.language].iloc[0]
    sales_val = model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'option_compare_sales') &
                                       (model_rep_main.vis_dict['variable'] == 'res_rub_distr'), model_rep_main.language].iloc[0]
    incr_sales = model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'option_compare_sales') &
                                        (model_rep_main.vis_dict['variable'] == 'res_rub_progress_distr'), model_rep_main.language].iloc[0]
    profit = model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'option_compare_sales') &
                                    (model_rep_main.vis_dict['variable'] == 'profit'), model_rep_main.language].iloc[0]
    romi = model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'option_compare_sales') &
                                  (model_rep_main.vis_dict['variable'] == 'romi'), model_rep_main.language].iloc[0]
    growth_vol = model_rep_main.vis_dict.loc[(model_rep_main.vis_dict['section'] == 'option_compare_sales') &
                                        (model_rep_main.vis_dict['variable'] == 'grw'),
                                        model_rep_main.language].iloc[0].replace("***", str(pred_year_start - 1))

    calc_df = pd.DataFrame(columns=[sales_col_name, sales_vol, sales_val, incr_sales, profit, romi, growth_vol])
    calc_df[sales_col_name] = list(options_df)[1:]

    ds, dt_pred, opt_sum = calc.option_report(model, ds, model_rep, options_m['zero'],
                                                opt_set, if_exact=if_exact)

    if isinstance(model_rep, list):
        opt_sum = opt_sum[-1]

    calc_df[calc_df[sales_col_name] == 'zero_short'] = ['zero_short',
                                                        opt_sum['pred_exact_0']['pred_exact_vol'],
                                                        opt_sum['pred_exact_0']['pred_exact_val'],
                                                        '-',
                                                        '-',
                                                        '-',
                                                        opt_sum['growth_vol'] - 1]
    calc_df[calc_df[sales_col_name] == 'zero_long'] = ['zero_long',
                                                       opt_sum['pred_long_vol'],
                                                       opt_sum['pred_long_val'],
                                                       '-',
                                                       '-',
                                                       '-',
                                                       opt_sum['pred_long_vol'] / opt_sum['ref_vol'] /
                                                       (pred_year_end - pred_year_start + 1) - 1]

    zero_val_short = opt_sum['pred_exact_0']['pred_exact_val']
    zero_val_long = opt_sum['pred_long_val']

    for option_name in options_names_subset:

        ds, dt_pred, opt_sum = calc.option_report(model, ds, model_rep, options_m[option_name],
                                                    opt_set, if_exact=if_exact)

        if isinstance(model_rep, list):
            opt_sum = opt_sum[-1]

        incr_sales_short = opt_sum['pred_exact_0']['pred_exact_val'] - zero_val_short
        profit_short = incr_sales_short - 1e6 * options_m[option_name]['bdg']
        romi_short = incr_sales_short / (1e6 * options_m[option_name]['bdg'])

        calc_df[calc_df[sales_col_name] == f'{option_name}_short'] = [f'{option_name}_short',
                                                                      opt_sum['pred_exact_0']['pred_exact_vol'],
                                                                      opt_sum['pred_exact_0']['pred_exact_val'],
                                                                      incr_sales_short,
                                                                      profit_short,
                                                                      romi_short,
                                                                      opt_sum['growth_vol'] - 1]

        incr_sales_long = opt_sum['pred_long_val'] - zero_val_long
        profit_long = incr_sales_long - 1e6 * options_m[option_name]['bdg'] * (pred_year_end - pred_year_start + 1)
        romi_long = incr_sales_long / (1e6 * options_m[option_name]['bdg'] * (pred_year_end - pred_year_start + 1))

        calc_df[calc_df[sales_col_name] == f'{option_name}_long'] = [f'{option_name}_long',
                                                                     opt_sum['pred_long_vol'],
                                                                     opt_sum['pred_long_val'],
                                                                     incr_sales_long,
                                                                     profit_long,
                                                                     romi_long,
                                                                     opt_sum['pred_long_vol'] / opt_sum['ref_vol'] /
                                                                     (pred_year_end - pred_year_start + 1) - 1]
    calc_df = calc_df.T
    new_header = calc_df.iloc[0]
    calc_df = calc_df[1:]
    calc_df.columns = new_header

    calc_df = calc_df.reset_index().rename(columns={'index': sales_col_name})

    calc_df[calc_df[sales_col_name] == growth_vol] = [f"{x:.2%}" if isinstance(x, (float, np.float64)) else x
                                                      for x in calc_df[calc_df[sales_col_name] == growth_vol].iloc[0]]

    calc_df[calc_df[sales_col_name] == romi] = [f"{round(x, 2)}" if isinstance(x, (float, np.float64)) else x
                                                for x in calc_df[calc_df[sales_col_name] == romi].iloc[0]]

    apply_cols = list(calc_df)[1:]
    for col in calc_df[apply_cols]:
        calc_df[col] = ['{:,.0f}'.format(x).replace(',', ' ') if isinstance(x, (float, np.float64))
                        else x for x in calc_df[col]]

    return options_df, calc_df


def _fill_table(table,
                ds: pd.DataFrame,
                pred_year_start: int,
                pred_year_end: int,
                last_row_bold: bool = False):
    """
    Fill table for options PPTX export with zebra coloration both horizontal and vertical (headers only).

    :param table: empty table (Table from python_pptx package)
    :param ds: dataset to be placed into `table`
    :param pred_year_start: first year with no / incomplete observed data
    :param pred_year_end: final year of prediction (mostly final year of dataset too)
    :param last_row_bold: highlight last row with bold font or not
    :return: filled table (Table from python_pptx package)
    """

    for n, x in enumerate(ds.columns):
        if n == 0:
            table.cell(0, n).text = x
        elif n % 2 == 1:
            table.cell(0, n).text = str(pred_year_start)
        else:
            table.cell(0, n).text = f"{str(pred_year_start)}-{str(pred_year_end)}"

        if n % 4 in [1, 2]:
            table.cell(0, n).fill.solid()
            table.cell(0, n).fill.fore_color.rgb = RGBColor(127, 127, 127)

    if last_row_bold:
        for index, row in ds.iterrows():
            for n, item in enumerate(row):
                if isinstance(item, str):
                    table.cell(index + 1, n).text = item
                else:
                    table.cell(index + 1, n).text = str(int(item))
                if index == len(ds) - 1:
                    table.cell(index + 1, n).text_frame.paragraphs[0].font.bold = True
    else:
        for index, row in ds.iterrows():
            for n, item in enumerate(row):
                if isinstance(item, str):
                    table.cell(index + 1, n).text = item
                else:
                    table.cell(index + 1, n).text = str(int(item))
    return table


